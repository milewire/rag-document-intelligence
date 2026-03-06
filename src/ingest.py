import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation
import os

def extract_text(file_path: str) -> str:
    """Extract text from a document based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        return extract_pdf(file_path)
    elif ext == ".docx":
        return extract_word(file_path)
    elif ext == ".xlsx":
        return extract_excel(file_path)
    elif ext == ".pptx":
        return extract_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def extract_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text.strip()

def extract_word(file_path: str) -> str:
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs]).strip()

def extract_excel(file_path: str) -> str:
    wb = openpyxl.load_workbook(file_path)
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text += f"\nSheet: {sheet}\n"
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join([str(cell) for cell in row if cell is not None])
            if row_text:
                text += row_text + "\n"
    return text.strip()

def extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for i, slide in enumerate(prs.slides):
        text += f"\nSlide {i+1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list:
    """Split text into overlapping chunks for better retrieval."""
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        start += chunk_size - overlap
    return chunks